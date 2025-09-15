"""
PowerPoint Agent for creating presentations (outputs pptx file)
"""

import re
import random
import os
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from GroqLLM import GroqLLM
from config import Config


class PowerPointAgent:
    def __init__(self, llm: GroqLLM):
        self.llm = llm

    def generate_slide_content(self, topic: str, slide_number: int, total_slides: int, 
                             slide_type: str, slide_focus: str) -> Dict:
        """Generate actual content for slides using LLM"""
        
        if slide_type == "title":
            prompt = f"""Create a compelling presentation title and subtitle about "{topic}".
            Respond ONLY in this format (plain text, no Markdown):
            Title: [Short, engaging title, max 7 words]
            Subtitle: [Concise subtitle, max 12 words]
            Visual Idea: [Optional suggestion for cover visual]
            """
        else:
            prompt = f"""Create slide content about "{topic}" focusing on: {slide_focus}.
            Respond ONLY in this format (plain text, no Markdown, no bulletpoint):
            Slide Title: [Concise title, max 7 words]
            - [Short, punchy phrase, <7 words]
            - [Short, punchy phrase, <7 words]
            - [Short, punchy phrase, <7 words]
            """

        try:
            response = self.llm.generate(prompt, max_tokens=400, model="gemma2-9b-it")
            
            if slide_type == "title":
                return self._parse_title_response(response, topic)
            else:
                return self._parse_content_response(response, topic, slide_number)
                
        except Exception as e:
            print(f"Content generation error: {e}")
            return self._get_fallback_content(topic, slide_number, slide_type)

    def _parse_title_response(self, response: str, topic: str) -> Dict:
        try:
            title = re.search(r'Title:\s*(.+)', response)
            subtitle = re.search(r'Subtitle:\s*(.+)', response)
            visual = re.search(r'Visual Idea:\s*(.+)', response)

            return {
                "title": title.group(1).strip() if title else f"{topic}: Overview",
                "subtitle": subtitle.group(1).strip() if subtitle else "Key Insights",
                "visual": visual.group(1).strip() if visual else None
            }
        except Exception:
            return {"title": f"{topic}", "subtitle": "Comprehensive Analysis", "visual": None}

    def _clean_text(self, text: str) -> str:
        """Remove markdown/extra symbols"""
        text = re.sub(r'[*#`>\\-]+', '', text)
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def _parse_content_response(self, response: str, topic: str, slide_number: int) -> Dict:
        try:
            # Match variations of "Slide Title"
            title_match = re.search(r'(?:Slide Title|Title)\s*[:\-]?\s*(.+)', response, re.IGNORECASE)
            title = self._clean_text(title_match.group(1)) if title_match else f"{topic} - Section {slide_number}"

            # Bullets (look for lines starting with - or •)
            bullets = re.findall(r'^[\-•]\s*(.+)', response, re.MULTILINE)
            bullets = [self._clean_text(b) for b in bullets if b.strip()]

            # Fallback if empty
            if not bullets:
                lines = [l for l in response.splitlines() if l.strip() and not l.lower().startswith("slide title")]
                bullets = [self._clean_text(l.lstrip("-• ").strip()) for l in lines][:3]

            # Clean leading dashes/stars just in case
            bullets = [re.sub(r'^[\-\*•]+\s*', '', b).strip() for b in bullets]

            while len(bullets) < 3:
                bullets.append(f"{topic} insight")

            return {"title": title, "bullets": bullets[:3], "visual": None}

        except Exception:
            return {
                "title": f"{topic} - Section {slide_number}",
                "bullets": [f"{topic} insight A", f"{topic} insight B", f"{topic} insight C"],
                "visual": None
            }

    def _get_fallback_content(self, topic: str, slide_number: int, slide_type: str) -> Dict:
        # TODO: instead of this maybe next try parsing through pypandoc
        if slide_type == "title":
            return {"title": topic, "subtitle": "Overview", "visual": None}
        else:
            return {"title": f"{topic} Slide", "bullets": ["Point A", "Point B", "Point C"], "visual": None}

    def _pick_theme(self) -> Dict:
        """Select a random color theme"""
        palettes = [
            {"bg": RGBColor(240, 248, 255), "accent": RGBColor(0, 102, 204)},   # Blue
            {"bg": RGBColor(240, 255, 240), "accent": RGBColor(34, 139, 34)},   # Green
            {"bg": RGBColor(255, 240, 240), "accent": RGBColor(178, 34, 34)},   # Red
        ]
        return random.choice(palettes)

    def _generate_outline(self, topic: str, total_slides: int) -> List[str]:
        """Generate an outline for the presentation"""
        prompt = f"""Create an outline for a {total_slides}-slide presentation about "{topic}".
        Provide a numbered list of slide focuses, each with a concise 3-5 word focus phrase.
        Do not repeat focuses. Avoid generic labels like 'introduction' or 'conclusion' more than once."""

        try:
            response = self.llm.generate(prompt, max_tokens=400, model="gemma2-9b-it")
            focuses = re.findall(r'\d+\. (.+)', response)
            return focuses[:total_slides] if focuses else [f"Aspect {i}" for i in range(1, total_slides+1)]
        except:
            return [f"Aspect {i}" for i in range(1, total_slides+1)]

    def create_presentation(self, topic: str = "Demo Topic", slides: int = 4) -> Dict:
        """Create a PowerPoint presentation"""
        try:
            prs = Presentation()
            theme = self._pick_theme()

            # Generate all focuses first
            focuses = self._generate_outline(topic, slides)

            # Title slide
            title_content = self.generate_slide_content(topic, 1, slides, "title", "")
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = title_content["title"]
            title_slide.placeholders[1].text = f"{title_content['subtitle']}\nAI-Generated Content"

            bg = title_slide.background.fill
            bg.solid()
            bg.fore_color.rgb = theme["bg"]

            # Content slides
            for i in range(2, slides + 1):
                slide_focus = focuses[i-1] if i-1 < len(focuses) else f"Section {i-1}"
                content = self.generate_slide_content(topic, i, slides, "content", slide_focus=slide_focus)

                slide = prs.slides.add_slide(prs.slide_layouts[1])

                # Background color
                bg = slide.background.fill
                bg.solid()
                bg.fore_color.rgb = theme["bg"]

                # Title formatting
                title_shape = slide.shapes.title
                title_shape.text = content["title"]
                title_shape.text_frame.paragraphs[0].font.size = Pt(32)
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.color.rgb = theme["accent"]

                # Bullets
                body_shape = slide.placeholders[1]
                body_shape.text = "\n".join([f"{b}" for b in content["bullets"]])

            # Save file
            safe_topic = re.sub(r'[^\w\s-]', '', topic)[:30]
            filename = f"AI_{safe_topic.replace(' ', '_')}.pptx"
            filepath = os.path.join(Config.OUTPUT_DIR, filename)
            
            prs.save(filepath)
            
            return {
                "success": True,
                "message": f"Created {slides}-slide deck on {topic}",
                "filename": filename,
                "filepath": filepath,
                "slides_count": slides,
                "topic": topic
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": f"PowerPoint creation failed: {str(e)}"
            }
        
    def create_presentation_from_content(self, processed_content: str, approach: str, slides: int, source_files: List[str], query: Optional[str] = None) -> Dict:
        """Create a PowerPoint from extracted file content (query-aware)"""
        try:
            prs = Presentation()
            theme = self._pick_theme()

            print(f"Query in create_presentation_from_content: {query}")

            # Generate everything at once
            all_slides = self._generate_slides_from_content(processed_content, slides, query=query)

            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = all_slides["titles"][0]

            # Build subtitle from first slide's bullets
            subtitle_text = "; ".join(all_slides["bulletpoints"][0])

            title_slide.placeholders[1].text = subtitle_text

            bg = title_slide.background.fill
            bg.solid()
            bg.fore_color.rgb = theme["bg"]

            # Remaining slides
            for i in range(1, len(all_slides["slides"])):
                slide = prs.slides.add_slide(prs.slide_layouts[1])

                bg = slide.background.fill
                bg.solid()
                bg.fore_color.rgb = theme["bg"]

                title_shape = slide.shapes.title
                title_shape.text = all_slides["titles"][i]
                title_shape.text_frame.paragraphs[0].font.size = Pt(32)
                title_shape.text_frame.paragraphs[0].font.bold = True
                title_shape.text_frame.paragraphs[0].font.color.rgb = theme["accent"]

                body_shape = slide.placeholders[1]
                body_shape.text = "\n".join(all_slides["bulletpoints"][i])

            # Save file
            safe_name = "Content_Presentation"
            filename = f"AI_{safe_name}.pptx"
            filepath = os.path.join(Config.OUTPUT_DIR, filename)
            prs.save(filepath)


            return {
                "success": True,
                "message": f"Created {slides}-slide deck from content",
                "filename": filename,
                "filepath": filepath,
                "slides_count": slides,
                "approach": approach,
                "query": query
            }

        except Exception as e:
            return {"success": False, "error": f"Content-based PowerPoint failed: {str(e)}"}
        
    def _generate_slides_from_content(self, content: str, total_slides: int, query: Optional[str] = None) -> List[Dict]:
        """Generate all slides in one call using simple structured format"""
        
        # Build focused instruction
        if query:
            focus_instruction = f"Focus specifically on '{query}' from the content below. Extract information related to {query}."
        else:
            focus_instruction = "Extract key information from the content below."
        
        print(f"Focus instruction: {focus_instruction}")

        # Truncate content to fit within token limits while preserving key information
        max_content_length = 8000  # Leave more room for the response
        if len(content) > max_content_length:
            # Try to keep query-relevant content if possible
            if query and query.lower() in content.lower():
                # Find sections that mention the query
                query_pos = content.lower().find(query.lower())
                start_pos = max(0, query_pos - 1000)
                content = content[start_pos:start_pos + max_content_length] + "..."
            else:
                content = content[:max_content_length] + "..."

        if total_slides is None:
            total_slides = 4  # Default to 4 slides if not specified

        print(f"We are creating a total of {total_slides} slides.")
        
        prompt = f"""{focus_instruction}

    Create exactly {total_slides} slides using the following structured format. Fill in everything between square brackets, keeping the rest of the template as is:

**SLIDE 1: [Main Topic Title]**
* [Key overview point 1]
* [Key overview point 2]

**SLIDE 2: [Specific Aspect]**
* [Key point 1]  
* [Key point 2]
* [Key point n]

**SLIDE [n]: [Another Key Aspect]**
* [Key point 1]
* [Key point 2]
* [Key point n]

- Follow the format exactly, for up to {total_slides} slides.
- Do not add extra slides, sections, text, or markdown.
- At least 3 key points per slide, max 5. There should be variation in the amount of key points per slide.
- Every point should be a concise, punchy phrase (max 10 words, no period at the end of the point).
- Titles should be engaging and informative (max 7 words).

    CONTENT TO ANALYZE:
    {content}"""

        try:
            response = self.llm.generate(prompt, max_tokens=1200, model="gemma2-9b-it")
            print(f"Structured response received, parsing...")
            print(f"Response: {response}")
            
            # Parse the structured text response
            slides = self._parse_slides(response)
            print(f"Parsed slides: {slides}")

            return slides
            
        except Exception as e:
            print(f"Slide generation error: {e}")
            return self._create_fallback_slides(total_slides, query, content)

    def _parse_slides(self, text: str) -> dict:
        # Regex to capture slide headers
        slide_pattern = re.compile(r"\*\*SLIDE\s+(\d+):\s*(.*?)\*\*", re.IGNORECASE)

        # Find all slide matches with their positions
        matches = list(slide_pattern.finditer(text))
        
        slides = []
        titles = []
        bulletpoints = []

        for i, match in enumerate(matches):
            slide_num = int(match.group(1))
            title = match.group(2).strip()

            # Start of this slide content
            start = match.end()
            # End is either the start of the next slide or end of text
            end = matches[i+1].start() if i + 1 < len(matches) else len(text)

            # Extract content between this slide and the next
            slide_content = text[start:end]

            # Collect bullet points (lines starting with *)
            bullets = [
                line.strip().lstrip("*").strip()
                for line in slide_content.splitlines()
                if line.strip().startswith("*")
            ]

            slides.append(slide_num)
            titles.append(title)
            bulletpoints.append(bullets)

        return {
            "slides": slides,
            "titles": titles,
            "bulletpoints": bulletpoints
        }
